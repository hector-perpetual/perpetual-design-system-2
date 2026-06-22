#!/usr/bin/env python3
"""
Genera el template de presentacion Perpetual Technologies (.pptx, 16:9).
Reinterpreta el layout "Essentials" con el design system Perpetual:
colores de marca, tipografia Armin Grotesk, logos oficiales y reglas duras
(nada de negro puro, fondos oscuros en #0b1220, logo correcto por contraste).

Fuente de verdad de tokens: perpetual-design-system (SKILL.md / tokens.md).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
LOGODIR = os.path.join(HERE, "perpetual-design-system", "assets", "logo")
LOGO_COLOR = os.path.join(LOGODIR, "perpetual-color.png")
LOGO_DARK = os.path.join(LOGODIR, "perpetual-dark.png")
OUT = os.path.join(HERE, "dist", "perpetual-deck-template.pptx")

# ---------------------------------------------------------------------------
# Tokens (de tokens.md)
# ---------------------------------------------------------------------------
def C(hexstr):
    return RGBColor.from_string(hexstr)

BRAND_BLUE = "0032CB"
BRAND_ORANGE = "F33D1F"
BRAND_YELLOW = "FBB900"

BG = "FFFFFF"
SURFACE = "F8F9FC"
SURFACE2 = "EEF1F8"
BORDER = "DDE1EF"
ACCENT = "1A56DB"
ACCENT2 = "F97316"
TEXT = "111827"
TEXT_DIM = "374151"
TEXT_MUTED = "6B7280"

BG_DARK = "0B1220"
SURFACE_DARK = "131C2E"
BORDER_DARK = "243047"
TEXT_ON_DARK = "FFFFFF"
TEXT_DIM_DARK = "9AA6BD"

OK = "059669"
WARN = "D97706"
ERROR = "DC2626"
VIOLET = "7E22CE"

# Paleta de datos para charts (orden recomendado en tokens.md)
DATA_PALETTE = [ACCENT, ACCENT2, OK, BRAND_YELLOW, VIOLET, TEXT_MUTED]

# Tipografia Armin Grotesk (instalada en el sistema; familias confirmadas por fc-list)
F_DISPLAY = "Armin Grotesk Black"     # H1 / titulos / numeros grandes (peso maximo con nombre de familia estable)
F_HEAD = "Armin Grotesk SemiBold"     # H2 / subtitulos de seccion
F_LABEL = "Armin Grotesk SemiBold"    # labels / chips uppercase
F_BODY = "Armin Grotesk"              # cuerpo
F_LIGHT = "Armin Grotesk Normal"      # subtitulos ligeros

YEAR = "2026"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(bg)
    return s


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _soft_shadow(shape, blur=12, dist=3, alpha=88):
    """Sombra suave estilo --shadow-card."""
    sp = shape._element.spPr
    # remove existing effect list
    for tag in ("a:effectLst",):
        el = sp.find(qn(tag))
        if el is not None:
            sp.remove(el)
    eff = sp.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(Emu(Pt(blur)).emu if False else int(blur * 12700)),
        "dist": str(int(dist * 12700)),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A2B5C"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha) * 1000))})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    sp.append(eff)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line)
        sp.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        _soft_shadow(sp)
    return sp


def grad_bar(s, x, y, w, h):
    """Barra/elemento con gradiente de marca azul->naranja."""
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sp)
    sp.line.fill.background()
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    # gradiente manual via XML
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, col in ((0, ACCENT), (50000, "3B82F6"), (100000, ACCENT2)):
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(pos)})
        c = gs.makeelement(qn("a:srgbClr"), {"val": col})
        gs.append(c)
        gsLst.append(gs)
    grad.append(gsLst)
    lin = grad.makeelement(qn("a:lin"), {"ang": "0", "scaled": "1"})
    grad.append(lin)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.insert(list(spPr).index(ln), grad)
    else:
        spPr.append(grad)
    return sp


def text(s, x, y, w, h, runs, size=18, color=TEXT, font=F_BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None,
         upper=False, wrap=True):
    """runs: str  OR  list de (texto, color, font, size_opt) para runs mixtos."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, color, font, size)]
    for item in runs:
        txt = item[0]
        rcolor = item[1] if len(item) > 1 and item[1] else color
        rfont = item[2] if len(item) > 2 and item[2] else font
        rsize = item[3] if len(item) > 3 and item[3] else size
        r = p.add_run()
        r.text = txt.upper() if upper else txt
        r.font.name = rfont
        r.font.size = Pt(rsize)
        r.font.bold = bold
        r.font.color.rgb = C(rcolor)
        if spacing is not None:
            _letter_spacing(r, spacing)
    return tb


def _letter_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def chip(s, x, y, label, fill=SURFACE2, fg=ACCENT, w=None):
    """Badge/chip uppercase tipo seccion."""
    w = w or (0.16 + 0.085 * len(label))
    sp = rect(s, x, y, w, 0.34, fill=fill, radius=0.5)
    text(s, x, y, w, 0.34, label, size=10.5, color=fg, font=F_LABEL,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2, upper=True)
    return w


def hexagon(s, x, y, size, fill=ACCENT, line=None):
    sp = rect(s, x, y, size, size, fill=fill, line=line, shape=MSO_SHAPE.HEXAGON)
    return sp


def logo(s, x, y, w, dark=False):
    path = LOGO_DARK if dark else LOGO_COLOR
    ratio = 287.0 / 2400.0
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w),
                               height=Inches(w * ratio))
    return pic


def footer(s, dark=False, page=None):
    fg = TEXT_DIM_DARK if dark else TEXT_MUTED
    logo(s, 0.55, 7.04, 1.05, dark=dark)
    text(s, 1.85, 7.0, 7, 0.3, "Confidencial  ·  Perpetual Technologies © " + YEAR,
         size=8.5, color=fg, font="JetBrains Mono", anchor=MSO_ANCHOR.MIDDLE)
    if page is not None:
        text(s, 11.7, 7.0, 1.08, 0.3, str(page).zfill(2), size=8.5, color=fg,
             font="JetBrains Mono", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(s, label, title_runs, sub=None, page=None):
    """Header estandar de slide de contenido."""
    logo(s, 0.55, 0.5, 1.25)
    chip(s, 11.0, 0.52, "ESSENTIALS" if False else "PERPETUAL", fill=SURFACE2, fg=ACCENT, w=1.35)
    text(s, 0.55, 1.15, 9.5, 0.32, label, size=11, color=ACCENT, font=F_LABEL,
         spacing=1.4, upper=True)
    text(s, 0.5, 1.42, 11.2, 0.9, title_runs, size=30, font=F_DISPLAY, color=TEXT)
    if sub:
        text(s, 0.55, 2.28, 9.5, 0.5, sub, size=13, color=TEXT_MUTED, font=F_LIGHT)


def line_chart(s, x, y, w, h, cats, series, smooth=True, colors=None, dark=False):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    colors = colors or DATA_PALETTE
    for i, ser in enumerate(plot.series):
        ser.smooth = smooth
        lf = ser.format.line
        lf.color.rgb = C(colors[i % len(colors)])
        lf.width = Pt(3.0)
    _style_axes(chart, dark)
    return chart


def col_chart(s, x, y, w, h, cats, series, colors=None, dark=False, gap=80, overlap=-20):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = gap
    try:
        plot.overlap = overlap
    except Exception:
        pass
    colors = colors or DATA_PALETTE
    for i, ser in enumerate(plot.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = C(colors[i % len(colors)])
        ser.format.line.fill.background()
    _style_axes(chart, dark)
    return chart


def donut(s, x, y, size, pct, color=ACCENT, rest=SURFACE2, hole=60):
    cd = CategoryChartData()
    cd.categories = ["v", "r"]
    cd.add_series("d", (pct, 100 - pct))
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(size), Inches(size), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    try:
        plot.donut_hole_size = hole
    except Exception:
        pass
    pts = plot.series[0].points
    pts[0].format.fill.solid()
    pts[0].format.fill.fore_color.rgb = C(color)
    pts[1].format.fill.solid()
    pts[1].format.fill.fore_color.rgb = C(rest)
    for p in pts:
        p.format.line.color.rgb = C(BG)
        p.format.line.width = Pt(1.5)
    return chart


def _style_axes(chart, dark=False):
    axc = TEXT_DIM_DARK if dark else TEXT_MUTED
    try:
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(9)
        ca.tick_labels.font.name = F_BODY
        ca.tick_labels.font.color.rgb = C(axc)
        ca.format.line.color.rgb = C(BORDER_DARK if dark else BORDER)
        ca.major_tick_mark = XL_TICK_MARK.NONE
        ca.minor_tick_mark = XL_TICK_MARK.NONE
        ca.has_major_gridlines = False
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.visible = False
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = C(BORDER_DARK if dark else BORDER)
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.color.rgb = C(axc)
        va.major_tick_mark = XL_TICK_MARK.NONE
    except Exception:
        pass


def card(s, x, y, w, h, fill=BG, radius=0.06, shadow=True, line=None):
    return rect(s, x, y, w, h, fill=fill, radius=radius, shadow=shadow, line=line)


def float_card(s, x, y, label, sub=None, w=2.85, icon=ACCENT, dark_icon=False,
               fg=TEXT, sub_fg=TEXT_MUTED, fill=BG):
    h = 0.92
    card(s, x, y, w, h, fill=fill, radius=0.16, shadow=True)
    hexagon(s, x + 0.18, y + 0.2, 0.52, fill=icon)
    rect(s, x + 0.30, y + 0.32, 0.28, 0.28, fill=BG if not dark_icon else SURFACE,
         shape=MSO_SHAPE.OVAL)
    tx = x + 0.92
    if sub:
        text(s, tx, y + 0.16, w - 1.0, 0.3, label, size=12, color=fg, font=F_HEAD)
        text(s, tx, y + 0.48, w - 1.0, 0.3, sub, size=10, color=sub_fg, font=F_BODY)
    else:
        text(s, tx, y, w - 1.0, h, label, size=12.5, color=fg, font=F_HEAD,
             anchor=MSO_ANCHOR.MIDDLE)


def stat(s, x, y, w, number, label, num_color=TEXT, lbl_color=TEXT_MUTED, num_size=34, align=PP_ALIGN.LEFT):
    text(s, x, y, w, 0.6, number, size=num_size, color=num_color, font=F_DISPLAY, align=align)
    text(s, x, y + num_size / 58.0, w, 0.35, label, size=11, color=lbl_color, font=F_LABEL,
         align=align, upper=True, spacing=0.6)


# ===========================================================================
# 01 — Portada / hero con line chart y cards flotantes
# ===========================================================================
def s01():
    s = slide(BG)
    rect(s, 0, 0, 13.333, 7.5, fill=SURFACE)         # base sutil
    rect(s, 0, 0, 13.333, 3.3, fill=BG, radius=0)
    logo(s, 0.55, 0.5, 1.45)
    chip(s, 11.45, 0.55, "TEMPLATE", fill=SURFACE2, fg=ACCENT, w=1.35)
    text(s, 1.0, 1.05, 11.3, 1.4,
         [("Resultados que ", TEXT, F_DISPLAY, 41), ("escalan", ACCENT, F_DISPLAY, 41),
          (" tu negocio.", TEXT, F_DISPLAY, 41)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    text(s, 2.5, 2.45, 8.3, 0.5,
         "Un vistazo a lo que construimos para marcas que quieren crecer con datos.",
         size=14, color=TEXT_MUTED, font=F_LIGHT, align=PP_ALIGN.CENTER)
    line_chart(s, 1.0, 3.15, 11.3, 3.7, ["Sem 1", "Sem 2", "Sem 3", "Sem 4", "Sem 5", "Sem 6"],
               [("Trafico", [210, 130, 150, 250, 300, 335])], colors=[ACCENT])
    float_card(s, 2.7, 3.35, "Trafico organico", w=2.7, icon=ACCENT)
    float_card(s, 7.7, 4.25, "Conversiones", w=2.7, icon=BG_DARK)
    footer(s, page=1)


# ===========================================================================
# 02 — Statement
# ===========================================================================
def s02():
    s = slide(BG)
    rect(s, -1.2, -1.2, 3.6, 3.6, fill=SURFACE2, shape=MSO_SHAPE.OVAL)
    rect(s, 10.9, 4.8, 3.8, 3.8, fill=SURFACE, shape=MSO_SHAPE.OVAL)
    rect(s, 11.6, 0.6, 1.3, 1.3, fill=SURFACE2, shape=MSO_SHAPE.OVAL)
    logo(s, 0.55, 0.5, 1.3)
    rect(s, 6.27, 2.35, 0.8, 0.07, fill=ACCENT, radius=0.5)
    text(s, 1.8, 2.7, 9.7, 2.0,
         [("Transformamos tu operacion en una ventaja ", TEXT, F_DISPLAY, 33),
          ("digital", ACCENT, F_DISPLAY, 33), (" medible.", TEXT, F_DISPLAY, 33)],
         align=PP_ALIGN.CENTER, line_spacing=1.05)
    footer(s, page=2)


# ===========================================================================
# 03 — Market growth (barras manuales + cards $)
# ===========================================================================
def s03():
    s = slide(BG)
    header(s, "Crecimiento", [("Crecimiento de mercado ", TEXT), (YEAR, ACCENT)])
    base = 6.3
    data = [("Q1", 126, 1.6), ("Q2", 256, 2.7), ("Q3", 345, 3.5)]
    bx = 0.9
    for i, (q, val, hh) in enumerate(data):
        x = bx + i * 1.75
        grad_bar(s, x, base - hh, 1.15, hh)
        text(s, x - 0.3, base + 0.05, 1.75, 0.3, q, size=12, color=TEXT_MUTED,
             font=F_LABEL, align=PP_ALIGN.CENTER, upper=True)
        card(s, x - 0.25, base - hh - 0.7, 1.6, 0.6, fill=BG, radius=0.18, shadow=True)
        text(s, x - 0.25, base - hh - 0.7, 1.6, 0.3, "S/ " + str(val) + "K", size=14,
             color=ACCENT, font=F_DISPLAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x - 0.25, base - hh - 0.42, 1.6, 0.25, "Ingresos", size=8.5,
             color=TEXT_MUTED, font=F_LABEL, align=PP_ALIGN.CENTER, upper=True)
    text(s, 7.1, 2.7, 5.4, 2.5,
         [("Cada trimestre crecimos con la misma logica: ", TEXT_DIM, F_BODY, 15),
          ("medir, ajustar y volver a medir.", TEXT, F_HEAD, 15)], line_spacing=1.3)
    text(s, 7.1, 4.3, 5.4, 2.0,
         "El ingreso recurrente paso de seis cifras bajas a un pipeline predecible "
         "en tres trimestres, sin inflar el gasto en medios.",
         size=12.5, color=TEXT_MUTED, font=F_BODY, line_spacing=1.35)
    footer(s, page=3)


# ===========================================================================
# 04 — Timeline
# ===========================================================================
def s04():
    s = slide(BG)
    header(s, "Trayectoria", [("Cada ano, un nuevo ", TEXT), ("hito.", ACCENT)])
    items = [("2021", "Primeros clientes", "Validamos el modelo con marcas locales."),
             ("2023", "Expansion regional", "Operacion en tres mercados de LatAm."),
             ("2025", "Producto propio", "Lanzamos la plataforma de medicion."),
             ("2026", "Best seller", "El servicio mas contratado del portafolio.")]
    lx = 5.2
    rect(s, lx, 2.7, 0.035, 3.9, fill=BORDER, radius=0)
    for i, (yr, t, d) in enumerate(items):
        y = 2.7 + i * 1.02
        hexagon(s, lx - 0.16, y, 0.36, fill=ACCENT if i % 2 == 0 else ACCENT2)
        text(s, lx - 1.7, y + 0.02, 1.45, 0.35, yr, size=15, color=ACCENT, font=F_DISPLAY,
             align=PP_ALIGN.RIGHT)
        text(s, lx + 0.5, y - 0.05, 6.8, 0.35, t, size=15, color=TEXT, font=F_HEAD)
        text(s, lx + 0.5, y + 0.3, 6.8, 0.35, d, size=11.5, color=TEXT_MUTED, font=F_BODY)
    footer(s, page=4)


# ===========================================================================
# 05 — Brand new / gear
# ===========================================================================
def s05():
    s = slide(BG)
    header(s, "A medida", [("Automatizacion con ", TEXT), ("esencia", ACCENT), (" a medida.", TEXT)])
    rect(s, 1.3, 3.0, 3.4, 3.4, fill=ACCENT, shape=MSO_SHAPE.GEAR_9)
    rect(s, 3.6, 4.6, 2.0, 2.0, fill=BG_DARK, shape=MSO_SHAPE.GEAR_6)
    rect(s, 2.55, 4.05, 1.0, 1.0, fill=BG, shape=MSO_SHAPE.OVAL)
    float_card(s, 7.4, 3.1, "Opcion 01", sub="4.314.000 procesos/ano", w=3.6, icon=ACCENT)
    float_card(s, 7.4, 4.45, "Opcion 02", sub="4.314.000 procesos/ano", w=3.6, icon=ACCENT2)
    text(s, 7.4, 5.85, 5.0, 0.6, "Flujos disenados para tu operacion, no plantillas genericas.",
         size=12, color=TEXT_MUTED, font=F_BODY, line_spacing=1.3)
    footer(s, page=5)


# ===========================================================================
# 06 — Bars + bubbles
# ===========================================================================
def s06():
    s = slide(BG)
    header(s, "Comunidad", [("Resultados que la gente ", TEXT), ("celebra.", ACCENT)])
    col_chart(s, 0.9, 3.0, 7.0, 3.7, ["Ene", "Feb", "Mar", "Abr", "May"],
              [("a", [40, 70, 55, 95, 80])],
              colors=[ACCENT], gap=60)
    bubbles = [(8.4, 3.1, "10K", "Alcance", ACCENT2),
               (10.2, 3.6, "22K", "Interaccion", ACCENT),
               (9.3, 5.0, "3K", "Leads", BG_DARK)]
    for x, y, n, lbl, col in bubbles:
        rect(s, x, y, 1.5, 1.5, fill=col, shape=MSO_SHAPE.OVAL, shadow=True)
        text(s, x, y + 0.32, 1.5, 0.5, n, size=22, color=TEXT_ON_DARK, font=F_DISPLAY,
             align=PP_ALIGN.CENTER)
        text(s, x, y + 0.85, 1.5, 0.3, lbl, size=9.5, color=TEXT_ON_DARK, font=F_LABEL,
             align=PP_ALIGN.CENTER, upper=True)
    footer(s, page=6)


# ===========================================================================
# 07 — Break section (dark)
# ===========================================================================
def s07():
    s = slide(BG_DARK)
    hexagon(s, 5.45, 0.7, 2.45, fill=SURFACE_DARK)
    hexagon(s, 6.25, 1.35, 0.85, fill=ACCENT)
    hexagon(s, 4.2, 5.3, 0.5, fill=ACCENT2)
    hexagon(s, 8.7, 5.0, 0.4, fill=BRAND_YELLOW)
    text(s, 1.5, 3.5, 10.3, 1.2,
         [("Seccion ", TEXT_ON_DARK, F_DISPLAY, 46), ("02", ACCENT, F_DISPLAY, 46)],
         align=PP_ALIGN.CENTER)
    text(s, 1.5, 4.7, 10.3, 0.5, "Estrategia y posicionamiento", size=15,
         color=TEXT_DIM_DARK, font=F_LIGHT, align=PP_ALIGN.CENTER)
    footer(s, dark=True, page=7)


# ===========================================================================
# 08 — World maps / cobertura
# ===========================================================================
def s08():
    s = slide(BG)
    header(s, "Cobertura", [("Presencia en la ", TEXT), ("region.", ACCENT)])
    card(s, 0.9, 2.9, 7.0, 3.5, fill=SURFACE, radius=0.05, shadow=False, line=BORDER)
    pins = [(2.2, 3.7), (3.6, 4.5), (4.8, 3.9), (5.9, 4.9), (6.7, 4.2)]
    for px, py in pins:
        rect(s, px, py, 0.45, 0.45, fill=ACCENT, shape=MSO_SHAPE.OVAL, shadow=True)
        rect(s, px + 0.13, py + 0.13, 0.19, 0.19, fill=BG, shape=MSO_SHAPE.OVAL)
    cards = [("Peru", "50"), ("Mexico", "67"), ("Colombia", "42"), ("Chile", "76")]
    for i, (pais, n) in enumerate(cards):
        x = 8.3 + (i % 2) * 2.35
        y = 2.9 + (i // 2) * 1.8
        card(s, x, y, 2.1, 1.55, fill=BG, radius=0.1, shadow=True)
        hexagon(s, x + 0.22, y + 0.25, 0.45, fill=DATA_PALETTE[i])
        text(s, x + 0.2, y + 0.78, 1.7, 0.5, n, size=26, color=TEXT, font=F_DISPLAY)
        text(s, x + 0.22, y + 1.2, 1.7, 0.3, pais, size=10.5, color=TEXT_MUTED,
             font=F_LABEL, upper=True)
    footer(s, page=8)


# ===========================================================================
# 09 — Donuts
# ===========================================================================
def s09():
    s = slide(BG)
    header(s, "Indicadores", [("Metricas que ", TEXT), ("importan.", ACCENT)])
    donut(s, 4.9, 3.1, 3.2, 38, color=ACCENT, rest=SURFACE2, hole=62)
    text(s, 4.9, 4.25, 3.2, 0.7, "38%", size=30, color=TEXT, font=F_DISPLAY,
         align=PP_ALIGN.CENTER)
    text(s, 4.9, 4.95, 3.2, 0.3, "Tasa de cierre", size=10, color=TEXT_MUTED,
         font=F_LABEL, align=PP_ALIGN.CENTER, upper=True)
    float_card(s, 1.0, 3.4, "Awesome!", sub="Nuestra meta", w=2.5, icon=ACCENT2)
    stat(s, 9.0, 3.2, 3.2, "350+", "Clientes activos", num_color=ACCENT)
    stat(s, 9.0, 4.6, 3.2, "550+", "Proyectos entregados", num_color=TEXT)
    footer(s, page=9)


# ===========================================================================
# 10 — Creative charts
# ===========================================================================
def s10():
    s = slide(BG)
    header(s, "Proyeccion", [("Proyeccion de ", TEXT), ("usuarios.", ACCENT)])
    cats = ["2024", "2025", "2026", "2028", "2030", "2032"]
    col_chart(s, 0.9, 2.9, 7.6, 3.7, cats,
              [("a", [30, 45, 95, 60, 80, 70])],
              colors=[ACCENT], gap=70)
    arrows = [(9.0, 3.0, "3.021+", "Usuarios nuevos", MSO_SHAPE.UP_ARROW, OK),
              (9.0, 4.9, "2.021+", "Usuarios recurrentes", MSO_SHAPE.DOWN_ARROW, ACCENT2)]
    for x, y, n, lbl, shp, col in arrows:
        rect(s, x, y, 0.5, 0.7, fill=col, shape=shp)
        text(s, x + 0.75, y - 0.05, 3.4, 0.5, n, size=22, color=TEXT, font=F_DISPLAY)
        text(s, x + 0.77, y + 0.42, 3.4, 0.3, lbl, size=10.5, color=TEXT_MUTED, font=F_BODY)
    footer(s, page=10)


# ===========================================================================
# 11 — Social stats
# ===========================================================================
def s11():
    s = slide(BG)
    header(s, "Redes", [("Tu marca en ", TEXT), ("redes.", ACCENT)])
    cats = ["Tweets", "Posts", "Likes", "Vistas"]
    vals = [156, 256, 556, 356]
    cols = [ACCENT, BG_DARK, ACCENT2, ACCENT]
    base = 6.2
    for i, (c, v) in enumerate(zip(cats, vals)):
        x = 1.4 + i * 2.5
        hh = v / 556.0 * 3.0
        rect(s, x, base - hh, 1.5, hh, fill=cols[i], radius=0.12)
        text(s, x, base + 0.1, 1.5, 0.35, str(v) + ",2", size=16, color=TEXT,
             font=F_DISPLAY, align=PP_ALIGN.CENTER)
        text(s, x, base + 0.5, 1.5, 0.3, c, size=10, color=TEXT_MUTED, font=F_LABEL,
             align=PP_ALIGN.CENTER, upper=True)
    text(s, 0.9, 6.95, 5, 0.3, "Hola! Eres increible.", size=11, color=ACCENT2, font=F_HEAD)
    footer(s, page=11)


# ===========================================================================
# 12 — Matrix / honeycomb
# ===========================================================================
def s12():
    s = slide(BG)
    header(s, "Metodo", [("Investigamos y creamos ", TEXT), ("ideas.", ACCENT)])
    text(s, 0.6, 3.2, 4.0, 2.0,
         "Un proceso en celdas: cada hipotesis se prueba antes de escalar.",
         size=14, color=TEXT_DIM, font=F_BODY, line_spacing=1.35)
    text(s, 0.6, 4.6, 4.0, 0.3, "Equipo Perpetual", size=11, color=TEXT_MUTED, font=F_LABEL, upper=True)
    cx, cy, hs = 8.0, 4.0, 1.35
    around = [(0, -1.15), (1.0, -0.58), (1.0, 0.58), (0, 1.15), (-1.0, 0.58), (-1.0, -0.58)]
    for dx, dy in around:
        hexagon(s, cx + dx * hs, cy + dy * hs, hs, fill=SURFACE2)
    hexagon(s, cx, cy, hs, fill=ACCENT)
    text(s, cx, cy + hs * 0.32, hs, 0.5, "Matrix", size=13, color=TEXT_ON_DARK,
         font=F_HEAD, align=PP_ALIGN.CENTER)
    footer(s, page=12)


# ===========================================================================
# 13 — Office charts (12 meses + donut)
# ===========================================================================
def s13():
    s = slide(BG)
    header(s, "Desempeno", [("Desempeno ", TEXT), ("mensual.", ACCENT)])
    meses = ["E", "F", "M", "A", "M", "J", "J", "A", "S", "O", "N", "D"]
    col_chart(s, 0.9, 2.9, 8.4, 3.4, meses,
              [("a", [50, 62, 45, 70, 66, 80, 74, 88, 79, 92, 85, 96])],
              colors=[ACCENT], gap=45)
    donut(s, 9.9, 3.4, 2.2, 63, color=ACCENT2, rest=SURFACE2, hole=60)
    text(s, 9.9, 4.2, 2.2, 0.5, "63%", size=22, color=TEXT, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    text(s, 9.7, 5.8, 2.7, 0.9, "Crecimiento sostenido mes a mes durante todo el ano.",
         size=11, color=TEXT_MUTED, font=F_BODY, align=PP_ALIGN.CENTER, line_spacing=1.3)
    footer(s, page=13)


# ===========================================================================
# 14 — Pricing
# ===========================================================================
def s14():
    s = slide(BG)
    header(s, "Planes", [("Planes que escalan ", TEXT), ("contigo.", ACCENT)])
    plans = [("Personal", "S/ 0", ["1 proyecto", "Reportes basicos", "Soporte por correo"], False),
             ("Equipo", "S/ 499", ["10 proyectos", "Dashboards en vivo", "Soporte prioritario"], True),
             ("Empresa", "S/ 1.299", ["Ilimitado", "Integraciones a medida", "Account manager"], False)]
    for i, (name, price, feats, hi) in enumerate(plans):
        x = 1.15 + i * 3.75
        y = 2.9
        fill = ACCENT if hi else BG
        fg = TEXT_ON_DARK if hi else TEXT
        muted = "DBE4FF" if hi else TEXT_MUTED
        card(s, x, y, 3.35, 3.7, fill=fill, radius=0.07, shadow=True,
             line=None if hi else BORDER)
        text(s, x + 0.35, y + 0.35, 2.6, 0.4, name, size=15, color=fg, font=F_HEAD, upper=False)
        text(s, x + 0.35, y + 0.8, 2.8, 0.6, price, size=30, color=fg, font=F_DISPLAY)
        text(s, x + 0.35, y + 1.45, 2.8, 0.3, "por mes", size=10, color=muted, font=F_LABEL, upper=True)
        for j, f in enumerate(feats):
            fy = y + 1.95 + j * 0.42
            rect(s, x + 0.35, fy + 0.06, 0.14, 0.14, fill=ACCENT2 if not hi else BRAND_YELLOW,
                 shape=MSO_SHAPE.OVAL)
            text(s, x + 0.62, fy, 2.5, 0.3, f, size=11.5, color=fg if hi else TEXT_DIM, font=F_BODY)
        btn_fill = BG if hi else ACCENT
        btn_fg = ACCENT if hi else TEXT_ON_DARK
        rect(s, x + 0.35, y + 3.15, 2.65, 0.42, fill=btn_fill, radius=0.5)
        text(s, x + 0.35, y + 3.15, 2.65, 0.42, "Elegir plan", size=11, color=btn_fg,
             font=F_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page=14)


# ===========================================================================
# 15 — Team
# ===========================================================================
def s15():
    s = slide(BG)
    header(s, "Equipo", [("Tres personas se suman al ", TEXT), ("equipo.", ACCENT)])
    people = [("Jazmin Galindo", "Marketing Lead", ACCENT, TEXT_ON_DARK,
               "Lidera la estrategia de adquisicion y medicion."),
              ("Kyron Downes", "Product Designer", BG_DARK, TEXT_ON_DARK,
               "Disena la experiencia de cada entregable."),
              ("Mara Ibanez", "Data Analyst", SURFACE, TEXT,
               "Convierte datos crudos en decisiones.")]
    for i, (name, role, fill, fg, quote) in enumerate(people):
        x = 1.0 + i * 3.85
        y = 3.0
        card(s, x, y, 3.45, 3.4, fill=fill, radius=0.07, shadow=True,
             line=BORDER if fill == SURFACE else None)
        hexagon(s, x + 0.35, y + 0.35, 0.75, fill=ACCENT2 if fill != ACCENT else BRAND_YELLOW)
        text(s, x + 1.3, y + 0.4, 2.0, 0.35, name, size=14, color=fg, font=F_HEAD)
        text(s, x + 1.3, y + 0.78, 2.0, 0.3, role, size=10.5,
             color=("DBE4FF" if fill == ACCENT else (TEXT_DIM_DARK if fill == BG_DARK else TEXT_MUTED)),
             font=F_LABEL, upper=True)
        text(s, x + 0.35, y + 1.5, 2.7, 0.3, "★ ★ ★ ★ ★", size=13,
             color=BRAND_YELLOW, font=F_BODY)
        text(s, x + 0.35, y + 2.05, 2.75, 1.0, quote, size=11.5,
             color=(fg if fill != SURFACE else TEXT_DIM), font=F_BODY, line_spacing=1.3)
    footer(s, page=15)


# ===========================================================================
# 16 — Roadmap
# ===========================================================================
def s16():
    s = slide(BG)
    header(s, "Roadmap", [("Hoja de ruta del ", TEXT), ("proyecto.", ACCENT)])
    y = 4.5
    rect(s, 1.1, y, 11.1, 0.12, fill=SURFACE2, radius=0.5)
    steps = [("01", "Descubrimiento", "Auditoria y objetivos."),
             ("02", "Estrategia", "Plan y metricas clave."),
             ("03", "Ejecucion", "Implementacion por sprints."),
             ("04", "Optimizacion", "Mejora continua con data.")]
    for i, (n, t, d) in enumerate(steps):
        x = 1.5 + i * 2.85
        up = i % 2 == 0
        hexagon(s, x, y - 0.27, 0.66, fill=DATA_PALETTE[i])
        text(s, x, y - 0.27, 0.66, 0.66, n, size=13, color=TEXT_ON_DARK, font=F_DISPLAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cy = (y - 1.65) if up else (y + 0.75)
        card(s, x - 0.55, cy, 2.5, 1.05, fill=BG, radius=0.12, shadow=True)
        text(s, x - 0.4, cy + 0.13, 2.25, 0.3, t, size=12.5, color=TEXT, font=F_HEAD)
        text(s, x - 0.4, cy + 0.5, 2.25, 0.45, d, size=10, color=TEXT_MUTED, font=F_BODY,
             line_spacing=1.2)
    footer(s, page=16)


# ===========================================================================
# 17 — New mission
# ===========================================================================
def s17():
    s = slide(BG)
    text(s, 0.4, 0.9, 8, 2.5, YEAR, size=150, color=SURFACE2, font=F_DISPLAY)
    header(s, "Mision", [("Nueva mision para ", TEXT), (YEAR, ACCENT)])
    cards = [("Puntos de medicion", "2,455,778", "Target anual", ACCENT),
             ("Misiones de empresa", "2,455,778", "Target anual", ACCENT2)]
    for i, (lbl, num, sub, col) in enumerate(cards):
        x = 1.0 + i * 5.9
        y = 4.0
        card(s, x, y, 5.3, 2.2, fill=SURFACE, radius=0.07, shadow=False, line=BORDER)
        hexagon(s, x + 0.4, y + 0.45, 0.7, fill=col)
        text(s, x + 1.45, y + 0.4, 3.6, 0.3, lbl, size=12, color=TEXT_MUTED, font=F_LABEL, upper=True)
        text(s, x + 1.45, y + 0.75, 3.7, 0.7, num, size=34, color=TEXT, font=F_DISPLAY)
        text(s, x + 1.45, y + 1.5, 3.7, 0.3, sub, size=11, color=TEXT_MUTED, font=F_BODY)
    text(s, 1.0, 6.45, 6, 0.3, "Presentacion por Equipo Perpetual", size=11,
         color=TEXT_MUTED, font=F_BODY)
    footer(s, page=17)


# ===========================================================================
# 18 — Branding 101
# ===========================================================================
def s18():
    s = slide(BG)
    logo(s, 0.55, 0.5, 1.3)
    text(s, 1.5, 1.4, 10.3, 0.8,
         [("Branding ", TEXT, F_DISPLAY, 30), ("101", ACCENT, F_DISPLAY, 30),
          (": que es la marca?", TEXT, F_DISPLAY, 30)],
         align=PP_ALIGN.CENTER)
    rect(s, 6.27, 2.25, 0.8, 0.07, fill=ACCENT, radius=0.5)
    card(s, 1.2, 2.9, 5.3, 1.6, fill=ACCENT, radius=0.1, shadow=True)
    text(s, 1.55, 3.15, 4.7, 0.35, "OPINION PRINCIPAL", size=10.5, color="DBE4FF",
         font=F_LABEL, upper=True, spacing=1.0)
    text(s, 1.55, 3.5, 4.8, 0.9,
         "La marca es la promesa que cumples cada vez que alguien te elige.",
         size=15, color=TEXT_ON_DARK, font=F_HEAD, line_spacing=1.2)
    rect(s, 6.9, 2.9, 2.5, 1.6, fill=SURFACE2, radius=0.16, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 9.6, 2.9, 2.5, 1.6, fill=SURFACE, radius=0.16, line=BORDER)
    card(s, 6.9, 4.7, 5.2, 1.6, fill=BG_DARK, radius=0.1, shadow=True)
    text(s, 7.25, 4.95, 4.5, 0.35, "TAREA ESPECIAL", size=10.5, color=TEXT_DIM_DARK,
         font=F_LABEL, upper=True, spacing=1.0)
    text(s, 7.25, 5.3, 4.6, 0.9, "Define tu promesa de marca para " + YEAR + " en una sola frase.",
         size=14, color=TEXT_ON_DARK, font=F_HEAD, line_spacing=1.2)
    rect(s, 1.2, 4.7, 5.3, 1.6, fill=SURFACE, radius=0.1, line=BORDER)
    text(s, 1.55, 5.0, 4.7, 1.1,
         "Una marca fuerte reduce el costo de adquisicion y sostiene el precio.",
         size=13, color=TEXT_DIM, font=F_BODY, line_spacing=1.3)
    footer(s, page=18)


# ===========================================================================
# 19 — 3D blocks (6 pasos)
# ===========================================================================
def s19():
    s = slide(BG)
    header(s, "Infografia", [("Infografia 3D: ", TEXT), ("6 pasos.", ACCENT)])
    blocks = [("DATA 01", ACCENT), ("DATA 02", "3B82F6"), ("DATA 03", BG_DARK),
              ("DATA 04", TEXT_MUTED), ("DATA 05", ACCENT2), ("DATA 06", BRAND_YELLOW)]
    bx, by = 2.4, 2.9
    for i, (lbl, col) in enumerate(blocks):
        y = by + i * 0.62
        x = bx - i * 0.12
        sp = rect(s, x, y, 2.6, 0.7, fill=col, shape=MSO_SHAPE.PARALLELOGRAM, radius=0)
        try:
            sp.adjustments[0] = 0.4
        except Exception:
            pass
        side = (i % 2 == 0)
        tx = (x + 3.0) if side else (x - 4.0)
        al = PP_ALIGN.LEFT if side else PP_ALIGN.RIGHT
        text(s, tx, y - 0.02, 3.9, 0.3, lbl, size=13, color=TEXT, font=F_DISPLAY, align=al)
        text(s, tx, y + 0.32, 3.9, 0.3, "Metrica clave del paso " + str(i + 1) + ".",
             size=10, color=TEXT_MUTED, font=F_BODY, align=al)
    footer(s, page=19)


# ===========================================================================
# Build
# ===========================================================================
for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
           s11, s12, s13, s14, s15, s16, s17, s18, s19]:
    fn()

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("OK:", OUT, "| slides:", len(prs.slides._sldIdLst))

